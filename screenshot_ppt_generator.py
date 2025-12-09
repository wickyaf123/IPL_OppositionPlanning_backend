"""
Screenshot-based PowerPoint Generator using Playwright
Captures exact visual appearance of web app slides at 1920x1080 resolution
"""

import os
import tempfile
import asyncio
from typing import List, Dict, Any
from urllib.parse import quote
from playwright.async_api import async_playwright, TimeoutError as PlaywrightTimeoutError
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class ScreenshotPPTGenerator:
    """Generate PowerPoint presentations from web app screenshots"""
    
    def __init__(self, base_url: str = "http://localhost:3000"):
        """
        Initialize generator
        
        Args:
            base_url: Base URL of the React app
        """
        self.base_url = base_url
        self.screenshot_width = 1920
        self.screenshot_height = 1080
        self.device_scale_factor = 2  # For sharper images (retina quality)
        
        # Brand colors for text elements
        self.colors = {
            'background': RGBColor(15, 23, 42),      # #0F172A
            'title': RGBColor(255, 255, 255),        # #FFFFFF
            'highlight': RGBColor(16, 185, 129),     # #10B981 (teal)
            'grey': RGBColor(156, 163, 175)          # #9CA3AF
        }
    
    async def capture_slide(self, url: str, output_path: str, timeout: int = 30000) -> bool:
        """
        Capture a single slide screenshot using Playwright
        
        Args:
            url: Full URL to capture
            output_path: Path to save screenshot
            timeout: Maximum wait time in milliseconds
            
        Returns:
            True if successful, False otherwise
        """
        try:
            print(f"Capturing screenshot from: {url}")
            
            async with async_playwright() as p:
                # Launch browser in headless mode
                browser = await p.chromium.launch(
                    headless=True,
                    args=['--no-sandbox', '--disable-setuid-sandbox']
                )
                
                # Create page with exact viewport size
                page = await browser.new_page(
                    viewport={
                        "width": self.screenshot_width,
                        "height": self.screenshot_height
                    },
                    device_scale_factor=self.device_scale_factor
                )
                
                # Navigate to page
                print(f"  Loading page...")
                await page.goto(url, wait_until="networkidle", timeout=timeout)
                
                # Wait for React to render
                print(f"  Waiting for content to load...")
                await asyncio.sleep(2)
                
                # Wait for slide ready signal (set by React component)
                try:
                    await page.wait_for_function(
                        "window.slideReady === true", 
                        timeout=15000
                    )
                    print(f"  Slide ready signal received")
                except PlaywrightTimeoutError:
                    print(f"  Warning: Slide ready signal timeout, proceeding anyway")
                
                # Additional wait for charts to render
                await asyncio.sleep(1)
                
                # Screenshot the export-slide container
                print(f"  Taking screenshot...")
                element = page.locator("#export-slide")
                await element.screenshot(path=output_path)
                
                print(f"  ✓ Screenshot saved: {output_path}")
                
                await browser.close()
                return True
                
        except Exception as e:
            print(f"  ✗ Error capturing slide: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    async def create_presentation(self, slides_data: List[Dict[str, Any]]) -> Presentation:
        """
        Create PowerPoint presentation from slide data
        
        Args:
            slides_data: List of slide configurations
            
        Returns:
            python-pptx Presentation object
        """
        # Create temporary directory for screenshots
        temp_dir = tempfile.mkdtemp()
        screenshot_paths = []
        
        print(f"\n{'='*60}")
        print(f"Starting PPT Generation with {len(slides_data)} slides")
        print(f"{'='*60}\n")
        
        try:
            # Capture screenshots for each slide
            for idx, slide_data in enumerate(slides_data):
                print(f"\n[Slide {idx + 1}/{len(slides_data)}]")
                
                # Build URL based on slide type
                url = self._build_slide_url(slide_data)
                screenshot_path = os.path.join(temp_dir, f"slide_{idx:02d}.png")
                
                if await self.capture_slide(url, screenshot_path):
                    screenshot_paths.append({
                        'path': screenshot_path,
                        'data': slide_data,
                        'index': idx + 1
                    })
                else:
                    print(f"  Warning: Failed to capture slide {idx + 1}, skipping...")
            
            print(f"\n{'='*60}")
            print(f"Successfully captured {len(screenshot_paths)}/{len(slides_data)} slides")
            print(f"{'='*60}\n")
            
            # Create PowerPoint presentation
            print("Creating PowerPoint presentation...")
            prs = Presentation()
            
            # Set slide size to 16:9 (matches 1920x1080)
            prs.slide_width = Inches(13.333)   # 1920/144 DPI
            prs.slide_height = Inches(7.5)     # 1080/144 DPI
            
            blank_layout = prs.slide_layouts[6]  # Blank layout
            
            for screenshot_info in screenshot_paths:
                print(f"  Adding slide {screenshot_info['index']} to presentation...")
                slide = prs.slides.add_slide(blank_layout)
                
                # Add screenshot as full-width image
                # Width = slide width, height auto-scales to maintain aspect ratio
                slide.shapes.add_picture(
                    screenshot_info['path'],
                    0,  # left position
                    0,  # top position
                    width=prs.slide_width
                )
                
                # Add editable Analyst Comments text box
                # Position it at the bottom, overlaying the screenshot
                self._add_analyst_comments_overlay(slide, screenshot_info['data'])
            
            print(f"\n✓ PowerPoint presentation created successfully!")
            print(f"  Total slides: {len(screenshot_paths)}")
            
            return prs
            
        finally:
            # Cleanup temporary files
            print(f"\nCleaning up temporary files...")
            for screenshot_info in screenshot_paths:
                try:
                    os.remove(screenshot_info['path'])
                except Exception as e:
                    print(f"  Warning: Could not remove {screenshot_info['path']}: {e}")
            try:
                os.rmdir(temp_dir)
            except Exception as e:
                print(f"  Warning: Could not remove temp directory: {e}")
    
    def _build_slide_url(self, slide_data: Dict[str, Any]) -> str:
        """Build URL for export route based on slide data"""
        slide_type = slide_data.get('type')
        params = [f"type={slide_type}"]
        
        if slide_type == 'player':
            player_name = slide_data.get('player_name', '')
            opposition = slide_data.get('opposition', '')
            params.append(f"player={quote(player_name)}")
            if opposition:
                params.append(f"opposition={quote(opposition)}")
                
        elif slide_type == 'team' or slide_type == 'overbyover':
            team_name = slide_data.get('team_name', '')
            params.append(f"team={quote(team_name)}")
            
        elif slide_type == 'venue':
            venue_name = slide_data.get('venue_name', '')
            params.append(f"venue={quote(venue_name)}")
        
        query_string = "&".join(params)
        return f"{self.base_url}/export/slide?{query_string}"
    
    def _add_analyst_comments_overlay(self, slide, slide_data: Dict[str, Any]):
        """
        Add editable analyst comments text box as overlay at bottom of slide
        This allows users to edit comments directly in PowerPoint
        """
        # Position at bottom of slide with semi-transparent background
        left = Inches(0.5)
        top = Inches(6.2)
        width = Inches(12.333)
        height = Inches(1.0)
        
        # Add a shape for background (optional, for better visibility)
        # You can comment this out if you want transparent background
        from pptx.enum.shapes import MSO_SHAPE
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark background
        background.line.color.rgb = RGBColor(16, 185, 129)  # Teal border
        background.line.width = Pt(1)
        
        # Add text box on top of background
        comments_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = comments_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        
        # Title
        p = text_frame.paragraphs[0]
        p.text = "Analyst Comments:"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['highlight']
        p.font.bold = True
        
        # Editable area
        text_frame.add_paragraph()
        p2 = text_frame.paragraphs[1]
        p2.text = "Click here to add your strategic insights and recommendations..."
        p2.font.name = 'Arial'
        p2.font.size = Pt(11)
        p2.font.color.rgb = self.colors['grey']
        p2.font.italic = True
