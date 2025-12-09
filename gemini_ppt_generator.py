"""
Gemini-powered PowerPoint Generator for IPL Opposition Analysis
Uses Google Gemini API to generate intelligent slide content
"""

import os
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Any
import json


class GeminiPPTGenerator:
    """Generate PowerPoint presentations using Gemini AI"""
    
    def __init__(self, api_key: str):
        """Initialize Gemini client with API key"""
        self.client = genai.Client(api_key=api_key)
        self.model = "gemini-2.0-flash-exp"  # Using latest Gemini model
        
        # Brand colors
        self.colors = {
            'background': RGBColor(15, 23, 42),      # #0F172A
            'card': RGBColor(17, 24, 39),            # #111827
            'title': RGBColor(255, 255, 255),        # #FFFFFF
            'body_text': RGBColor(203, 213, 225),    # #CBD5E1
            'highlight': RGBColor(16, 185, 129),     # #10B981
            'warning': RGBColor(245, 158, 11),       # #F59E0B
            'grey': RGBColor(156, 163, 175)          # #9CA3AF
        }
    
    def generate_slide_content(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Use Gemini to generate intelligent slide content based on data
        
        Args:
            slide_data: Dictionary containing slide type, player/team name, and statistics
            
        Returns:
            Dictionary with generated content including insights, recommendations, etc.
        """
        prompt = self._create_prompt(slide_data)
        
        try:
            response = self.client.models.generate_content(
                model=self.model,
                contents=prompt
            )
            
            # Parse the response as JSON
            content = response.text
            # Clean markdown code blocks if present
            if content.startswith("```json"):
                content = content.replace("```json", "").replace("```", "").strip()
            elif content.startswith("```"):
                content = content.replace("```", "").strip()
                
            return json.loads(content)
        except Exception as e:
            print(f"Error generating content with Gemini: {e}")
            return self._get_fallback_content(slide_data)
    
    def _create_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create a detailed prompt for Gemini based on slide type"""
        slide_type = slide_data.get('type', 'unknown')
        
        if slide_type == 'player':
            return self._create_player_prompt(slide_data)
        elif slide_type == 'team':
            return self._create_team_prompt(slide_data)
        elif slide_type == 'overbyover':
            return self._create_overbyover_prompt(slide_data)
        elif slide_type == 'venue':
            return self._create_venue_prompt(slide_data)
        else:
            return self._create_generic_prompt(slide_data)
    
    def _create_player_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create prompt for player analysis slide"""
        player_name = slide_data.get('player_name', 'Unknown Player')
        stats = slide_data.get('statistics', {})
        
        return f"""You are a professional cricket analyst creating a PowerPoint slide for IPL opposition planning.

Player: {player_name}
Statistics: {json.dumps(stats, indent=2)}

Generate a comprehensive analysis in JSON format with the following structure:
{{
    "title": "Player Analysis: {player_name}",
    "key_insights": [
        "3-4 bullet points with key performance insights based on the statistics"
    ],
    "strengths": [
        "3-4 specific strengths with data-backed observations"
    ],
    "weaknesses": [
        "2-3 areas for improvement or vulnerabilities"
    ],
    "tactical_recommendations": [
        "3-4 specific tactical recommendations for bowling/fielding strategies"
    ],
    "summary": "A concise 2-3 sentence summary of the player's overall profile"
}}

Focus on:
- Strike rates against different bowling types
- Performance in different phases (powerplay, middle, death)
- Boundary hitting ability
- Dot ball percentages
- Recent form and trends

Return ONLY valid JSON, no additional text."""
    
    def _create_team_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create prompt for team analysis slide"""
        team_name = slide_data.get('team_name', 'Unknown Team')
        stats = slide_data.get('statistics', {})
        
        return f"""You are a professional cricket analyst creating a PowerPoint slide for IPL opposition planning.

Team: {team_name}
Statistics: {json.dumps(stats, indent=2)}

Generate a comprehensive team analysis in JSON format:
{{
    "title": "Opposition Analysis: {team_name}",
    "key_insights": [
        "4-5 bullet points about team's overall performance and patterns"
    ],
    "batting_strengths": [
        "3-4 points about batting lineup strengths"
    ],
    "bowling_patterns": [
        "3-4 points about bowling attack patterns"
    ],
    "tactical_recommendations": [
        "4-5 specific strategies to exploit weaknesses"
    ],
    "match_situation_analysis": "2-3 sentences about how the team performs in different match situations"
}}

Return ONLY valid JSON, no additional text."""
    
    def _create_overbyover_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create prompt for over-by-over analysis slide"""
        team_name = slide_data.get('team_name', 'Unknown Team')
        over_data = slide_data.get('over_by_over_data', {})
        
        return f"""You are a professional cricket analyst creating a PowerPoint slide for IPL opposition planning.

Team: {team_name}
Over-by-Over Bowling Data: {json.dumps(over_data, indent=2)}

Generate an over-by-over bowling analysis in JSON format:
{{
    "title": "Over-by-Over Analysis: {team_name}",
    "powerplay_analysis": [
        "3-4 points about overs 1-6 bowling patterns"
    ],
    "middle_overs_analysis": [
        "3-4 points about overs 7-15 bowling patterns"
    ],
    "death_overs_analysis": [
        "3-4 points about overs 16-20 bowling patterns"
    ],
    "bowler_rotation_insights": [
        "3-4 points about how the team rotates bowlers"
    ],
    "tactical_recommendations": [
        "4-5 specific recommendations for batting strategy in each phase"
    ]
}}

Return ONLY valid JSON, no additional text."""
    
    def _create_venue_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create prompt for venue analysis slide"""
        venue_name = slide_data.get('venue_name', 'Unknown Venue')
        stats = slide_data.get('statistics', {})
        
        return f"""You are a professional cricket analyst creating a PowerPoint slide for IPL opposition planning.

Venue: {venue_name}
Statistics: {json.dumps(stats, indent=2)}

Generate a comprehensive venue analysis in JSON format:
{{
    "title": "Venue Analysis: {venue_name}",
    "pitch_characteristics": [
        "3-4 points about pitch behavior and conditions"
    ],
    "historical_trends": [
        "3-4 points about historical performance at this venue"
    ],
    "toss_impact": [
        "2-3 points about toss decisions and their impact"
    ],
    "tactical_recommendations": [
        "4-5 specific recommendations for team selection and strategy"
    ],
    "weather_factors": "1-2 sentences about typical weather conditions and their impact"
}}

Return ONLY valid JSON, no additional text."""
    
    def _create_generic_prompt(self, slide_data: Dict[str, Any]) -> str:
        """Create generic prompt for unknown slide types"""
        return f"""Generate cricket analysis content in JSON format based on this data:
{json.dumps(slide_data, indent=2)}

Return structured insights in JSON format with title, key_insights, and recommendations."""
    
    def _get_fallback_content(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """Provide fallback content if Gemini API fails"""
        slide_type = slide_data.get('type', 'unknown')
        name = slide_data.get('player_name') or slide_data.get('team_name') or slide_data.get('venue_name', 'Unknown')
        
        return {
            "title": f"{slide_type.title()} Analysis: {name}",
            "key_insights": [
                f"Analysis for {name}",
                "Performance metrics and trends",
                "Strategic considerations",
                "Tactical recommendations"
            ],
            "tactical_recommendations": [
                "Analyze recent performance data",
                "Consider match situation and conditions",
                "Adapt strategy based on opposition",
                "Monitor player form and fitness"
            ],
            "summary": f"Comprehensive analysis for {name} with data-driven insights."
        }
    
    def create_presentation(self, slides_data: List[Dict[str, Any]]) -> Presentation:
        """
        Create a complete PowerPoint presentation with Gemini-generated content
        
        Args:
            slides_data: List of dictionaries containing slide information
            
        Returns:
            python-pptx Presentation object
        """
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        for idx, slide_data in enumerate(slides_data):
            print(f"Generating slide {idx + 1}/{len(slides_data)}...")
            
            # Generate content using Gemini
            content = self.generate_slide_content(slide_data)
            
            # Create slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.colors['background']
            
            # Add title
            self._add_title(slide, content.get('title', 'Analysis'))
            
            # Add content based on what's available
            self._add_slide_content(slide, content)
            
            # Add slide number
            self._add_slide_number(slide, idx + 1, len(slides_data))
        
        return prs
    
    def _add_title(self, slide, title_text: str):
        """Add title to slide"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        p = title_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(32)
        p.font.color.rgb = self.colors['title']
        p.font.bold = True
    
    def _add_slide_content(self, slide, content: Dict[str, Any]):
        """Add generated content to slide"""
        y_position = 1.2
        
        # Add key insights
        if 'key_insights' in content:
            y_position = self._add_section(slide, "Key Insights", content['key_insights'], y_position)
        
        # Add strengths (for player slides)
        if 'strengths' in content:
            y_position = self._add_section(slide, "Strengths", content['strengths'], y_position, self.colors['highlight'])
        
        # Add weaknesses (for player slides)
        if 'weaknesses' in content:
            y_position = self._add_section(slide, "Areas for Improvement", content['weaknesses'], y_position, self.colors['warning'])
        
        # Add tactical recommendations
        if 'tactical_recommendations' in content:
            self._add_section(slide, "Tactical Recommendations", content['tactical_recommendations'], y_position)
        
        # Add analyst comments box
        self._add_analyst_comments(slide)
    
    def _add_section(self, slide, title: str, bullets: List[str], y_pos: float, color=None) -> float:
        """Add a section with title and bullet points"""
        if color is None:
            color = self.colors['highlight']
        
        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(15), Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.color.rgb = color
        p.font.bold = True
        
        # Bullet points
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.5), Inches(14.5), Inches(1.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets[:4]):  # Limit to 4 bullets per section
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = f"• {bullet}"
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['body_text']
            p.space_after = Pt(8)
        
        return y_pos + 2.2
    
    def _add_analyst_comments(self, slide):
        """Add editable analyst comments section"""
        comments_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(15), Inches(1.5))
        comments_frame = comments_box.text_frame
        
        # Title
        comments_frame.text = "Analyst Comments:"
        p = comments_frame.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = self.colors['highlight']
        p.font.bold = True
        
        # Editable text area
        comments_frame.add_paragraph()
        p2 = comments_frame.paragraphs[1]
        p2.text = "Click here to add your strategic insights and recommendations..."
        p2.font.name = 'Arial'
        p2.font.size = Pt(12)
        p2.font.color.rgb = self.colors['grey']
        p2.font.italic = True
    
    def _add_slide_number(self, slide, current: int, total: int):
        """Add slide number"""
        number_box = slide.shapes.add_textbox(Inches(14), Inches(8.3), Inches(1.5), Inches(0.4))
        number_frame = number_box.text_frame
        number_frame.text = f"Slide {current} of {total}"
        p = number_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.color.rgb = self.colors['grey']
